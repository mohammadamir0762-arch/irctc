"""
Generates the project presentation (PS91_Presentation.pptx) in the repo root.

Every figure in the deck comes from either analysis/benchmark.py (the model
work) or new/EDA_INSIGHTS_REPORT.md (the delay study). Nothing is rounded up
or restated more strongly than the source supports — the delay slides carry
their n=100 sample size on the slide itself, because that is the first thing
a reviewer will ask about.

Run:  python analysis/make_ppt.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
CHARTS = ROOT / "new"
OUTPUT = ROOT / "PS91_Presentation.pptx"

# Palette — light ground so the white-background matplotlib charts sit on the
# page as figures rather than glowing panels; dark navy reserved for the title
# and section breaks so the deck has rhythm.
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
GOOD = RGBColor(0x05, 0x96, 0x69)
BAD = RGBColor(0xDC, 0x26, 0x26)
WARN = RGBColor(0xD9, 0x77, 0x06)

FONT = "Arial"
W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs, bg=PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False
    return slide


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def write(tf, text, size, color=INK, bold=False, space_after=6, first=False,
          align=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.space_after = Pt(space_after)
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = FONT
    return p


def rule(slide, x, y, w, color=ACCENT, thickness=Pt(3)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False


def heading(slide, title, eyebrow=None):
    y = Inches(0.5)
    if eyebrow:
        tf = textbox(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.3))
        write(tf, eyebrow.upper(), 11, MUTED, bold=True, first=True, space_after=0)
        y += Inches(0.34)
    tf = textbox(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.6))
    write(tf, title, 30, INK, bold=True, first=True, space_after=0)
    rule(slide, MARGIN, y + Inches(0.66), Inches(1.1))
    return y + Inches(1.02)


def title_slide(prs):
    slide = blank(prs, INK)
    tf = textbox(slide, MARGIN, Inches(2.15), W - 2 * MARGIN, Inches(0.4))
    write(tf, "PROBLEM STATEMENT PS91", 13, ACCENT, bold=True, first=True, space_after=0)

    tf = textbox(slide, MARGIN, Inches(2.72), Inches(10.4), Inches(1.9))
    write(tf, "Railway Ticket", 52, PAPER, bold=True, first=True, space_after=2)
    write(tf, "Confirmation Prediction", 52, PAPER, bold=True, space_after=0)

    tf = textbox(slide, MARGIN, Inches(4.75), Inches(9.2), Inches(0.9))
    write(tf, "Predicting whether a waitlisted Indian Railways ticket will confirm,",
          17, RGBColor(0xCB, 0xD5, 0xE1), first=True, space_after=2)
    write(tf, "using a model trained on 28,748 real waitlisted tickets.",
          17, RGBColor(0xCB, 0xD5, 0xE1), space_after=0)

    rule(slide, MARGIN, Inches(5.95), Inches(2.2), ACCENT)
    tf = textbox(slide, MARGIN, Inches(6.2), Inches(11), Inches(0.6))
    write(tf, "Live demo   irctc-smoky.vercel.app", 13, PAPER, first=True, space_after=3)
    write(tf, "Source   github.com/mohammadamir0762-arch/irctc", 13, MUTED, space_after=0)


def section_slide(prs, number, title, subtitle):
    slide = blank(prs, INK)
    tf = textbox(slide, MARGIN, Inches(2.6), Inches(2), Inches(1.2))
    write(tf, number, 64, RGBColor(0x1E, 0x29, 0x3B), bold=True, first=True, space_after=0)

    tf = textbox(slide, MARGIN, Inches(3.5), Inches(10.5), Inches(1.4))
    write(tf, title, 40, PAPER, bold=True, first=True, space_after=6)
    write(tf, subtitle, 16, RGBColor(0x94, 0xA3, 0xB8), space_after=0)
    rule(slide, MARGIN, Inches(3.35), Inches(1.6))


def bullets_slide(prs, title, bullets, eyebrow=None, footnote=None):
    slide = blank(prs)
    y = heading(slide, title, eyebrow)
    tf = textbox(slide, MARGIN, y, W - 2 * MARGIN, H - y - Inches(0.9))
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            label, detail = item
            p = write(tf, label, 19, INK, bold=True, first=(i == 0), space_after=2)
            write(tf, detail, 15, MUTED, space_after=16)
        else:
            write(tf, "•   " + item, 17, INK, first=(i == 0), space_after=12)
    if footnote:
        tf = textbox(slide, MARGIN, H - Inches(0.82), W - 2 * MARGIN, Inches(0.4))
        write(tf, footnote, 11, MUTED, first=True, space_after=0, italic=True)


def stat_row(slide, y, stats, box_h=Inches(1.5)):
    """Evenly spaced KPI tiles."""
    gap = Inches(0.26)
    total = W - 2 * MARGIN
    box_w = int((total - gap * (len(stats) - 1)) / len(stats))
    for i, (value, label, color) in enumerate(stats):
        x = MARGIN + i * (box_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WASH
        card.line.fill.background()
        card.shadow.inherit = False
        card.adjustments[0] = 0.06

        tf = textbox(slide, x, y + Inches(0.22), box_w, Inches(0.7), PP_ALIGN.CENTER)
        write(tf, value, 30, color, bold=True, first=True, space_after=0, align=PP_ALIGN.CENTER)
        tf = textbox(slide, x + Inches(0.1), y + Inches(0.93), box_w - Inches(0.2),
                     Inches(0.5), PP_ALIGN.CENTER)
        write(tf, label, 11, MUTED, first=True, space_after=0, align=PP_ALIGN.CENTER)


def table_slide(prs, title, headers, rows, eyebrow=None, widths=None, footnote=None,
                row_colors=None):
    slide = blank(prs)
    y = heading(slide, title, eyebrow)
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = min(Inches(0.46) * n_rows, H - y - Inches(1.0))
    shape = slide.shapes.add_table(n_rows, n_cols, MARGIN, y, W - 2 * MARGIN, height)
    table = shape.table
    table.first_row = True

    if widths:
        total = W - 2 * MARGIN
        for i, frac in enumerate(widths):
            table.columns[i].width = Emu(int(total * frac))

    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = PAPER
                run.font.name = FONT

    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if r % 2 else WASH
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            colour = INK
            if row_colors and row_colors.get(r - 1) and c == len(row) - 1:
                colour = row_colors[r - 1]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = colour
                    run.font.bold = bool(row_colors and row_colors.get(r - 1) and c == len(row) - 1)
                    run.font.name = FONT

    if footnote:
        tf = textbox(slide, MARGIN, H - Inches(0.8), W - 2 * MARGIN, Inches(0.4))
        write(tf, footnote, 11, MUTED, first=True, space_after=0, italic=True)
    return slide


def chart_slide(prs, title, image, takeaway, eyebrow=None, caveat=None):
    """Chart on the left, the point it makes on the right — so a reader gets
    the conclusion without decoding axes."""
    from PIL import Image

    slide = blank(prs)
    y = heading(slide, title, eyebrow)

    avail_w, avail_h = Inches(7.6), H - y - Inches(0.75)
    iw, ih = Image.open(image).size
    scale = min(avail_w / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(image), MARGIN + int((avail_w - w) / 2),
                             y + int((avail_h - h) / 2), w, h)

    tx = MARGIN + avail_w + Inches(0.4)
    tw = W - tx - MARGIN
    tf = textbox(slide, tx, y + Inches(0.1), tw, avail_h)
    write(tf, "WHAT IT SHOWS", 10, ACCENT, bold=True, first=True, space_after=8)
    for i, line in enumerate(takeaway):
        write(tf, line, 14, INK if i == 0 else MUTED, bold=(i == 0), space_after=10)
    if caveat:
        write(tf, caveat, 11, WARN, space_after=0, italic=True)


def build():
    prs = new_deck()

    # ---------------------------------------------------------------- intro
    title_slide(prs)

    bullets_slide(
        prs, "The problem", eyebrow="Problem statement PS91",
        bullets=[
            ("Waitlisted passengers are left guessing",
             "A ticket says \"WL 23\". Nothing tells you whether that will actually confirm, "
             "so people hold a ticket they may not be able to use, or cancel one that would have cleared."),
            ("The decision has a real deadline",
             "Chart preparation happens roughly 4 hours before departure. Cancel too early and you "
             "lose a seat that would have opened; cancel too late and you pay a higher penalty."),
            ("The information exists, but is not surfaced",
             "Waitlist position, quota, class and time remaining all influence the outcome. "
             "None of it is turned into an actual probability for the passenger."),
        ],
    )

    slide = blank(prs)
    y = heading(slide, "What we built", eyebrow="Solution")
    tf = textbox(slide, MARGIN, y, W - 2 * MARGIN, Inches(1.0))
    write(tf, "Enter a 10-digit PNR. The system looks up the ticket, predicts the chance it confirms, "
              "and explains which factors drove that number.", 17, INK, first=True, space_after=0)
    stat_row(slide, y + Inches(1.1), [
        ("28,748", "REAL TICKETS TRAINED ON", INK),
        ("0.796", "TEST AUC", ACCENT),
        ("4", "MODEL INPUTS", INK),
        ("2", "LIVE DEPLOYMENTS", INK),
    ])
    tf = textbox(slide, MARGIN, y + Inches(2.95), W - 2 * MARGIN, Inches(1.8))
    write(tf, "Web app  ·  irctc-smoky.vercel.app", 15, INK, bold=True, first=True, space_after=6)
    write(tf, "REST API  ·  pnr-predictor-api.onrender.com", 15, INK, bold=True, space_after=6)
    write(tf, "Mobile app  ·  React Native (Expo), builds cleanly, not yet device-tested",
          15, MUTED, space_after=6)
    write(tf, "Source  ·  github.com/mohammadamir0762-arch/irctc", 15, MUTED, space_after=0)

    bullets_slide(
        prs, "How it works", eyebrow="Pipeline",
        bullets=[
            ("1.  Look up the PNR",
             "A third-party IRCTC API returns train, class, quota, booking status and current status."),
            ("2.  Normalise the response",
             "Providers disagree on field names and status formats. One adapter converts every "
             "provider into a single internal shape, so swapping providers changes nothing downstream."),
            ("3.  Derive model inputs",
             "Travel class, waitlist position at booking, current waitlist position, days until journey. "
             "All read directly from the response — no estimated or invented values."),
            ("4.  Predict and explain",
             "A gradient boosting model returns a probability, plus the top factors behind it."),
        ],
    )

    # ----------------------------------------------------------------- data
    section_slide(prs, "01", "Getting real data",
                  "The hardest part of this project was not the model")

    bullets_slide(
        prs, "There is no official dataset", eyebrow="The core obstacle",
        bullets=[
            "IRCTC publishes no historical dataset of PNR outcomes, and offers no bulk data API.",
            "Apps like ConfirmTkt and Trainman are IRCTC Authorised Partners — a commercial "
            "arrangement, not a developer signup.",
            "Our first model was trained on data we generated ourselves. We scrapped it.",
            "It was circular: the labels came from a formula we wrote, so the model could only "
            "re-learn our own assumptions. It would have scored well and known nothing.",
        ],
        footnote="Deleting that model was the single most important decision in the project.",
    )

    table_slide(
        prs, "We audited three public datasets", eyebrow="Dataset selection",
        headers=["Dataset", "Rows", "Finding", "Verdict"],
        widths=[0.26, 0.11, 0.45, 0.18],
        rows=[
            ["Railway Waitinglist", "53,381",
             "Real. Class mix matches reality; confirmation falls cleanly with waitlist position.",
             "USED"],
            ["Railway Ticket Confirmation", "30,000",
             "Fabricated. Sequential PNRs, 30,000 unique dates for 30,000 rows, Shatabdi trains "
             "with sleeper class.", "REJECTED"],
            ["Railofy (competition)", "36,775",
             "Real and richer, but values are encoded. Cannot be fed live data.", "BENCHMARK ONLY"],
        ],
        row_colors={0: GOOD, 1: BAD, 2: WARN},
        footnote="Only one of the three could actually train a deployable model.",
    )

    bullets_slide(
        prs, "How we caught the fabricated dataset", eyebrow="Data validation",
        bullets=[
            ("PNR numbers were sequential",
             "PNR0000000000, PNR0000000001, PNR0000000002 — incrementing by exactly one."),
            ("30,000 rows, 30,000 unique journey dates",
             "One journey per calendar day, in perfect order. Real bookings cluster."),
            ("Impossible train and class combinations",
             "Shatabdi trains listed with Sleeper class, spread perfectly evenly across every "
             "class. Shatabdi is chair-car only."),
            ("The label was a tautology",
             "Every row with a waitlist position was \"Not Confirmed\", every row without was "
             "\"Confirmed\". A model scores a perfect AUC of 1.0 and learns nothing."),
        ],
        footnote="Training on this would have produced an impressive, meaningless number.",
    )

    # ---------------------------------------------------------------- model
    section_slide(prs, "02", "The model", "Four inputs, and the evidence that they are the right ones")

    slide = blank(prs)
    y = heading(slide, "Model and inputs", eyebrow="Approach")
    tf = textbox(slide, MARGIN, y, Inches(5.9), Inches(4.4))
    write(tf, "Algorithm", 13, ACCENT, bold=True, first=True, space_after=4)
    write(tf, "Histogram gradient boosting classifier", 16, INK, space_after=14)
    write(tf, "Training data", 13, ACCENT, bold=True, space_after=4)
    write(tf, "28,748 real waitlisted tickets, 39,724 observations", 16, INK, space_after=4)
    write(tf, "Each ticket is observed at 30, 7 and 2 days before travel, matching what the app "
              "knows at prediction time.", 13, MUTED, space_after=14)
    write(tf, "Validation", 13, ACCENT, bold=True, space_after=4)
    write(tf, "Grouped train/test split", 16, INK, space_after=4)
    write(tf, "Observations of the same ticket never appear on both sides of the split, so the "
              "score is not inflated by leakage.", 13, MUTED, space_after=0)

    tf = textbox(slide, MARGIN + Inches(6.4), y, Inches(5.1), Inches(4.4))
    write(tf, "The four inputs", 13, ACCENT, bold=True, first=True, space_after=8)
    for label in ["Travel class", "Waitlist position when booked",
                  "Current waitlist position", "Days until journey"]:
        write(tf, "•   " + label, 16, INK, space_after=8)
    write(tf, "Every one is read straight from the PNR lookup. Nothing is estimated, and nothing "
              "is used in training that the live system cannot obtain.", 13, MUTED, space_after=0)

    table_slide(
        prs, "Are four inputs actually enough?", eyebrow="Evidence from real tickets",
        headers=["Factor", "Measured effect", "Spread"],
        widths=[0.30, 0.50, 0.20],
        rows=[
            ["Waitlist position", "12.9% confirm at WL 1–5, down to 0.2% at WL 60+", "65×"],
            ["Days until journey", "6.6% at 30 days, 12.9% at 7 days, 15.1% at 2 days", "2.3×"],
            ["Travel class", "2A 5.2%, 3A 8.6%, SL 12.5%, CC 33.6%", "6.5×"],
        ],
        footnote="Measured on the real dataset, holding the other factors fixed. "
                 "These are the determining factors, not filler.",
    )

    slide = blank(prs)
    y = heading(slide, "How much signal do we capture?", eyebrow="Benchmark")
    tf = textbox(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.8))
    write(tf, "We benchmarked against a Kaggle competition dataset with 36,775 labelled tickets "
              "and 23 features, to find the ceiling.", 16, INK, first=True, space_after=0)
    stat_row(slide, y + Inches(0.85), [
        ("0.500", "RANDOM BASELINE", MUTED),
        ("0.796", "OUR MODEL  ·  4 FEATURES  ·  LIVE", ACCENT),
        ("0.945", "CEILING  ·  23 FEATURES  ·  NOT DEPLOYABLE", INK),
    ], box_h=Inches(1.6))
    tf = textbox(slide, MARGIN, y + Inches(2.75), W - 2 * MARGIN, Inches(1.6))
    write(tf, "We reach 66% of the achievable signal above random, using 4 inputs instead of 23.",
          20, INK, bold=True, first=True, space_after=10)
    write(tf, "Permutation importance across all 23 features ranks current waitlist position "
              "first, by more than double the next feature. Our model holds the strongest "
              "predictor there is.", 15, MUTED, space_after=0)

    bullets_slide(
        prs, "Why we did not ship the higher score", eyebrow="An honest trade-off",
        bullets=[
            ("The 0.945 model cannot accept a real PNR",
             "Railofy encodes waitlist position as a fraction — 1/2, 1/3, 3/8 — of a "
             "denominator that is not in the file. A real \"WL 25\" cannot be converted into it."),
            ("We tested the workaround rather than assuming",
             "Quantile-mapping real values onto its distribution and scoring against real labels "
             "gives AUC 0.480 — worse than a coin flip."),
            ("So it stays a benchmark",
             "Shipping it would have looked better on a slide and performed worse for every user. "
             "The deployed model uses raw values it can actually read."),
        ],
        footnote="analysis/benchmark.py in the repo reproduces every number on these slides.",
    )

    # -------------------------------------------------------------- product
    section_slide(prs, "03", "Engineering", "What runs, and what it runs on")

    bullets_slide(
        prs, "Architecture", eyebrow="Stack",
        bullets=[
            ("Backend  ·  Python, FastAPI, scikit-learn",
             "Six endpoints. /pnr for the main flow, /predict for direct model access, /model "
             "publishes the model's own provenance and scores so the numbers can be checked."),
            ("Frontend  ·  HTML, CSS, JavaScript",
             "No build step. Detects whether it is running locally or deployed and picks the "
             "right API automatically."),
            ("Mobile  ·  React Native via Expo",
             "Same flow as the web app, pointed at the same deployed API."),
            ("Hosting  ·  Vercel and Render, both free tiers",
             "Deployed from GitHub. Cold starts are handled with an explicit \"waking up\" "
             "message rather than a frozen screen."),
        ],
    )

    bullets_slide(
        prs, "Working within a 10-request limit", eyebrow="Constraint",
        bullets=[
            ("The free API tier allows 10 requests per month",
             "Not per day. We read this from the live rate-limit headers rather than the pricing "
             "page, which is behind a login."),
            ("So the live site runs on simulated data, and says so",
             "A banner marks every simulated result. The real integration is built and verified "
             "against a genuine API response — it is switched off, not missing."),
            ("Every check is logged for future training",
             "Quota, train and route are recorded even though the current model cannot use them. "
             "That is the path to the features we are missing."),
        ],
        footnote="A demo that runs out of quota mid-presentation is worse than an honest simulated one.",
    )

    # --------------------------------------------------------- delay study
    section_slide(prs, "04", "Railway delay analysis",
                  "A separate exploratory study on train punctuality")

    slide = blank(prs)
    y = heading(slide, "Scope of the delay study", eyebrow="Read this first")
    tf = textbox(slide, MARGIN, y, W - 2 * MARGIN, Inches(1.0))
    write(tf, "This is a separate exploratory analysis of train punctuality, not part of the "
              "prediction model. It is a small sample and the findings are indicative, not conclusive.",
          16, INK, first=True, space_after=0)
    stat_row(slide, y + Inches(1.0), [
        ("100", "JOURNEY RECORDS", INK),
        ("~10", "RECORDS PER TRAIN", WARN),
        ("2016–2025", "PERIOD COVERED", INK),
        ("Winter", "ONLY SEASON PRESENT", WARN),
    ])
    tf = textbox(slide, MARGIN, y + Inches(2.85), W - 2 * MARGIN, Inches(1.6))
    write(tf, "What this sample cannot support", 15, INK, bold=True, first=True, space_after=8)
    write(tf, "•   Any claim about Indian Railways nationally — 100 records is far too "
              "small, and the 6% on-time rate reflects this sample, not the network.", 14, MUTED, space_after=6)
    write(tf, "•   Seasonal conclusions — every record is from winter.", 14, MUTED, space_after=6)
    write(tf, "•   Year-on-year trends — roughly 10 records per year is too thin to "
              "separate a real trend from noise.", 14, MUTED, space_after=0)

    delay_charts = [
        ("Delay distribution", "01_delay_distribution.png",
         ["Most delays are small, a few are extreme.",
          "The distribution is right-skewed: the bulk of journeys sit under 50 minutes, "
          "with a long tail reaching 244 minutes.",
          "Median 25 min vs mean 36.3 min — the gap is caused by that tail, so the median "
          "is the fairer summary."], None),
        ("On-time performance", "04_ontime_performance.png",
         ["In this sample, 94 of 100 journeys arrived late.",
          "Only 6 journeys met the on-time threshold."],
         "This is a 100-record sample, not a national statistic."),
        ("Performance by train", "02_train_performance.png",
         ["A 36-minute gap separates best from worst.",
          "Vivek Express averages 20.4 min; Yesvantpur–Howrah averages 56.6 min.",
          "Variation this wide points to route and operational factors rather than "
          "something inherent to running trains."],
         "Around 10 records per train — indicative only."),
        ("Spread within each train", "09_train_boxplot.png",
         ["Consistency varies as much as the average.",
          "Some trains are reliably a little late; others swing between on-time and "
          "severely delayed.",
          "For a passenger, that predictability matters as much as the mean."], None),
        ("Distance vs delay", "03_distance_vs_delay.png",
         ["Distance alone does not predict delay.",
          "Correlation is −0.04 — essentially none.",
          "Grouped into bands, long routes do average higher (52.6 min vs 30.5 min), so the "
          "relationship is not linear and distance is a weak proxy for something else."],
         "The near-zero correlation and the band differences must be read together."),
        ("Delay by distance band", "08_distance_category.png",
         ["Long routes average the highest delay.",
          "Short 30.5 min, medium 36.5 min, long 52.6 min.",
          "More opportunities to accumulate delay, and more shared track."], None),
        ("Service frequency", "06_frequency_impact.png",
         ["More frequent services performed worse here.",
          "Daily 38.1 min, weekly 34.5 min, tri-weekly 27.4 min.",
          "Consistent with congestion and tighter turnarounds on busy corridors, though this "
          "sample cannot establish the cause."], None),
        ("Delay over time", "07_temporal_trend.png",
         ["Year-to-year figures move sharply.",
          "The sample shows a low around 2019–2020 and a peak in 2023.",
          "With roughly 10 records per year, these swings are not reliable evidence of a trend."],
         "Too few records per year to draw conclusions from."),
        ("Route comparison", "10_routes_comparison.png",
         ["The same distance can produce very different delays.",
          "Dibrugarh–Kanyakumari runs 4,198 km at 20.4 min average, better than several "
          "far shorter routes.",
          "This is the strongest signal in the study: the specific corridor matters more "
          "than the distance travelled."], None),
        ("Seasonal view", "05_seasonal_analysis.png",
         ["No seasonal comparison is possible.",
          "Every record in the dataset is from winter.",
          "Included to show the gap, not a finding. Collecting year-round data is the "
          "obvious next step."], None),
    ]
    for title, filename, takeaway, caveat in delay_charts:
        path = CHARTS / filename
        if path.exists():
            chart_slide(prs, title, path, takeaway, eyebrow="Delay analysis", caveat=caveat)

    bullets_slide(
        prs, "What the delay study suggests", eyebrow="Delay analysis  ·  summary",
        bullets=[
            "Route identity matters more than route length — the clearest pattern in the data.",
            "Delay is dominated by a minority of journeys; the median journey is far better than "
            "the mean suggests.",
            "Higher-frequency services performed worse in this sample, which is worth investigating "
            "properly with more data.",
            "The sample is too small and too seasonally narrow to support conclusions about the "
            "network as a whole.",
        ],
        footnote="Treated as an exploratory study that motivates further data collection.",
    )

    # -------------------------------------------------------------- closing
    section_slide(prs, "05", "Limitations and next steps",
                  "What we know does not work yet")

    bullets_slide(
        prs, "Known limitations", eyebrow="Being straight about it",
        bullets=[
            ("The model uses four features, not everything that matters",
             "Quota ranks fifth in importance and we would like it — General confirms 26%, "
             "Pooled 44%, Remote Location 45%. No public dataset provides it with real outcomes "
             "in a usable form."),
            ("Live data is limited by a 10-request monthly quota",
             "The integration is built and verified; it runs in simulated mode to stay within it."),
            ("The waitlisted path has not been seen against a live waitlisted PNR",
             "Confirmed tickets are verified end-to-end. The waitlist branch is derived from the "
             "provider schema."),
            ("The mobile app has not been run on a physical device",
             "It bundles without errors, but that is not the same as tested."),
        ],
    )

    bullets_slide(
        prs, "Next steps", eyebrow="Roadmap",
        bullets=[
            ("Collect our own outcome data",
             "Logging is already in place and records quota, train and route. With paid API quota "
             "this becomes a training set containing the features we are missing."),
            ("Retrain with quota, train and route",
             "The benchmark shows these are worth roughly 0.15 AUC — the gap between 0.796 and 0.945."),
            ("Benchmark against the provider's own prediction",
             "The API returns its own confidence figure. Logging it beside ours and the real "
             "outcome shows which is better calibrated."),
            ("Verify on device and extend the delay study",
             "Run the mobile app on real hardware; collect year-round delay data to make the "
             "seasonal analysis possible."),
        ],
    )

    slide = blank(prs, INK)
    tf = textbox(slide, MARGIN, Inches(2.35), Inches(11), Inches(1.0))
    write(tf, "Thank you", 44, PAPER, bold=True, first=True, space_after=0)
    rule(slide, MARGIN, Inches(3.5), Inches(1.6))
    tf = textbox(slide, MARGIN, Inches(3.85), Inches(11.5), Inches(2.4))
    write(tf, "Live demo   irctc-smoky.vercel.app", 17, PAPER, first=True, space_after=10)
    write(tf, "API   pnr-predictor-api.onrender.com", 17, RGBColor(0xCB, 0xD5, 0xE1), space_after=10)
    write(tf, "Source   github.com/mohammadamir0762-arch/irctc", 17, RGBColor(0xCB, 0xD5, 0xE1),
          space_after=22)
    write(tf, "Every figure in this deck is reproducible with  python analysis/benchmark.py",
          13, MUTED, space_after=0, italic=True)

    prs.save(OUTPUT)
    return len(prs.slides.__iter__.__self__._sldIdLst)


if __name__ == "__main__":
    count = build()
    print(f"Wrote {OUTPUT}")
    print(f"{count} slides")
