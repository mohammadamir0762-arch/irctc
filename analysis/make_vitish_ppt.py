"""
Fills the VITISH'26 / SIH idea template with this project's content.

Works on a copy of the supplied template rather than rebuilding a deck, so the
event branding, logos, footer bar and page numbers stay exactly as issued. Only
the body text of each section is replaced; the section headings the template
calls "idea details pointers" are preserved, and each one is answered under a
bold label of the same name.

The instructions slide is removed, as the template permits.

Run:  python analysis/make_vitish_ppt.py
"""

import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "VITISH'26_IdeaTemplate.pptx"
OUTPUT = ROOT / "PS91_VITISH26_Idea.pptx"

INK = RGBColor(0x1F, 0x29, 0x37)
LABEL = RGBColor(0x0B, 0x3D, 0x91)   # deep blue, sits with the template's navy
MUTED = RGBColor(0x44, 0x4C, 0x56)

BODY_FONT = "Arial"

# Anything the team must supply is marked with this so it cannot be missed.
TBD = "«fill in»"


def body_box(slide):
    """The template puts one free TEXT_BOX per slide for the answer content.
    Pick the largest non-placeholder text box on the slide."""
    boxes = [
        sh for sh in slide.shapes
        if sh.has_text_frame and not sh.is_placeholder and sh.width and sh.height
    ]
    if not boxes:
        return None
    return max(boxes, key=lambda s: s.width * s.height)


def set_body(slide, blocks, size=13, label_size=15, box=None, top=None, height=None,
             width=None):
    """Replace a slide's body text with labelled bullet blocks."""
    shape = box or body_box(slide)
    if shape is None:
        return
    if top is not None:
        shape.top = top
    if height is not None:
        shape.height = height
    if width is not None:
        shape.width = width

    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for label, points in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = label
        p.space_before = Pt(0 if not p == tf.paragraphs[0] else 0)
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.size = Pt(label_size)
            run.font.bold = True
            run.font.color.rgb = LABEL
            run.font.name = BODY_FONT

        for point in points:
            bp = tf.add_paragraph()
            bp.text = "•  " + point
            bp.space_after = Pt(3)
            for run in bp.runs:
                run.font.size = Pt(size)
                run.font.bold = False
                run.font.color.rgb = INK
                run.font.name = BODY_FONT

        spacer = tf.add_paragraph()
        spacer.text = ""
        for run in spacer.runs:
            run.font.size = Pt(5)


def set_placeholder(slide, contains, lines, size=28, bold=True, color=None):
    """Replace a title placeholder whose current text contains `contains`."""
    for sh in slide.shapes:
        if not (sh.has_text_frame and sh.text_frame.text.strip()):
            continue
        if contains.lower() not in sh.text_frame.text.lower():
            continue
        tf = sh.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.space_after = Pt(2)
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.name = "Times New Roman"
                if color is not None:
                    run.font.color.rgb = color
        return sh
    return None


def drop_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    rId = slides[index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)
    sldIdLst.remove(slides[index])


def build():
    shutil.copy(TEMPLATE, OUTPUT)
    prs = Presentation(OUTPUT)
    s = list(prs.slides)

    # ------------------------------------------------------------ 1  title
    for sh in s[0].shapes:
        if sh.has_text_frame and "Problem Statement ID" in sh.text_frame.text:
            tf = sh.text_frame
            tf.word_wrap = True
            tf.clear()
            rows = [
                ("Problem Statement ID:", "PS91"),
                ("Problem Statement Title:", "Railway Ticket Confirmation Prediction"),
                ("Theme:", "Transportation & Logistics"),
                ("PS Category:", "Software"),
                ("Team Name:", TBD),
            ]
            for i, (k, v) in enumerate(rows):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{k} {v}"
                p.space_after = Pt(8)
                for run in p.runs:
                    run.font.size = Pt(16)
                    run.font.bold = False
                    run.font.color.rgb = INK
                    run.font.name = BODY_FONT
            break
    set_placeholder(s[0], "TITLE PAGE",
                    ["RailSure", "Journey confidence for Indian Railways"],
                    size=30)

    # ------------------------------------------------------- 2  idea title
    set_placeholder(s[1], "IDEA TITLE",
                    ["RailSure — predicting whether a journey will actually work"],
                    size=26)
    set_body(s[1], [
        ("Proposed Solution", [
            "One platform answering a single question: will this journey work, end to end?",
            "Confirmation probability for any waitlisted PNR, with the factors behind it",
            "Route search returning direct trains and multi-leg connecting options",
            "Risk score for missing the second train when the first one runs late",
            "Control-room view of how many passengers transfer between two trains at a junction",
        ]),
        ("How it addresses the problem", [
            "Replaces \"will it confirm?\" guesswork with a probability the passenger can act on",
            "Treats a two-leg trip as one journey instead of two unrelated tickets",
            "Makes the passenger handover visible to staff while it can still be acted upon",
        ]),
        ("Innovation and uniqueness", [
            "Existing apps stop at PNR prediction; none score the connection itself",
            "No current system tells the railway who on a late train is about to miss a departure",
            "The same prediction serves both the passenger and the operator",
        ]),
    ], size=13, label_size=15, top=Inches(1.7), height=Inches(5.1))

    # -------------------------------------------------- 3  technical approach
    set_body(s[2], [
        ("Technologies to be used", [
            "Backend: Python, FastAPI, scikit-learn (gradient boosting classifier)",
            "Web: HTML, CSS, JavaScript  |  Mobile: React Native (Expo)",
            "Data: SQLite journey log; REST API consumed by both clients",
            "Deployment: Render (API) and Vercel (web), continuous deploy from GitHub",
        ]),
        ("Methodology and process for implementation", [
            "1.  PNR lookup — a provider-agnostic adapter normalises any vendor response",
            "2.  Prediction — the model scores each leg from class, waitlist position and days left",
            "3.  Route building — direct and connecting options generated between two cities",
            "4.  Connection risk — leg-one delay history compared against the layover window",
            "5.  Aggregation — journeys grouped by junction and train pair for the operations view",
        ]),
    ], size=13, label_size=15, top=Inches(2.3), height=Inches(4.4))

    # ------------------------------------------------ 4  feasibility/viability
    set_body(s[3], [
        ("Analysis of the feasibility of the idea", [
            "Working prototype already deployed and publicly reachable (web app and REST API)",
            "Prediction model trained on real waitlisted-ticket data and validated on held-out tickets",
            "Entire stack runs on free-tier infrastructure, so cost is not a barrier to piloting",
        ]),
        ("Potential challenges and risks", [
            "Indian Railways publishes no official bulk data API for ticket status",
            "Third-party API access is rate-limited, restricting live traffic at prototype stage",
            "Public datasets do not carry quota and route together with confirmation outcomes",
            "Operations view requires railway-side integration to reach its full value",
        ]),
        ("Strategies for overcoming these challenges", [
            "Provider-agnostic adapter — vendors can be swapped without touching application code",
            "Every lookup is logged, so the system accumulates its own training data through use",
            "IRCTC Authorised Partner route is the established path to production-scale access",
            "Passenger features deliver value on their own, before any railway integration",
        ]),
    ], size=13, label_size=15, top=Inches(2.2), height=Inches(4.6), width=Inches(12.0))

    # -------------------------------------------------------- 5  impact
    set_body(s[4], [
        ("Potential impact on the target audience", [
            "Passengers decide with a number rather than a guess before committing to a trip",
            "Connecting journeys can be planned with a realistic buffer instead of hope",
            "Staff can see how many people a short hold would protect, and decide accordingly",
        ]),
        ("Benefits of the solution", [
            "Social — fewer passengers stranded overnight at junction stations; elderly travellers, "
            "families and long-distance migrants benefit most",
            "Economic — fewer wasted fares, cancellation penalties and emergency rebookings",
            "Operational — a missed connection becomes a decision made in time rather than an "
            "incident discovered afterwards",
            "Systemic — aggregated transfer demand shows where timetables genuinely need adjusting",
        ]),
    ], size=13, label_size=15, top=Inches(2.2), height=Inches(4.6))

    # ---------------------------------------------------- 6  team capability
    set_body(s[5], [
        ("Team's experience in the domain", [
            f"{TBD} — members, roles and relevant coursework or internships",
        ]),
        ("Why choosing the domain", [
            "Waitlisted tickets and missed connections affect a very large number of travellers daily",
            "The problem is data-rich and measurable, which makes it well suited to a machine "
            "learning solution",
            "Every team member has personally experienced the problem being solved",
        ]),
        ("Proof of past work in the domain", [
            "Deployed prototype: irctc-smoky.vercel.app",
            "Public API: pnr-predictor-api.onrender.com",
            "Full source code and analysis: github.com/mohammadamir0762-arch/irctc",
            f"{TBD} — any earlier projects, hackathons or publications",
        ]),
    ], size=13, label_size=15, top=Inches(2.3), height=Inches(4.5))

    # ------------------------------------------------------- 7  references
    set_body(s[6], [
        ("Prototype and source", [
            "Live web application — irctc-smoky.vercel.app",
            "Public prediction API — pnr-predictor-api.onrender.com",
            "Source code, model training and analysis — github.com/mohammadamir0762-arch/irctc",
        ]),
        ("Data sources", [
            "Railway Waitinglist Dataset (Kaggle) — real waitlist movement and confirmation outcomes",
            "Indian Railways Challenge (Kaggle) — used to benchmark feature selection",
            "Indian Railways train timetable and route data — data.gov.in",
        ]),
        ("Domain references", [
            "CRIS — Centre for Railway Information Systems, operator of the reservation database",
            "IRCTC PNR enquiry and Authorised Partner programme",
            "IndianRailAPI and RapidAPI IRCTC endpoints — third-party PNR status access",
        ]),
    ], size=13, label_size=15, top=Inches(2.0), height=Inches(4.8), width=Inches(12.0))

    # remove the instructions slide (template explicitly allows this)
    if len(list(prs.slides)) > 7:
        drop_slide(prs, 7)

    prs.save(OUTPUT)
    return len(prs.slides._sldIdLst)


if __name__ == "__main__":
    print(f"Wrote {OUTPUT}\n{build()} slides")
